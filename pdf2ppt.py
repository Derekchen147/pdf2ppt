import fitz
import os
import time
from tqdm import tqdm

import os
import datetime
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt


for filename in os.listdir('source_files/'):
    # 将pdf一张张切割成jpg
    timestamp = time.time()
    print('切割pdf为jpg...')
    # base_path = './source_files'                  # 输入要转换的PDF所在的文件夹
    # filename = os.listdir(base_path)              # 获取PDF文件列表
    full_path = "./source_files/%s" % filename            # 拼接，得到PDF文件的绝对路径
    print(full_path)
    doc = fitz.open(full_path)  # 打开一个PDF文件，doc为Document类型，是一个包含每一页PDF文件的列表
    rotate = int(0)  # 设置图片的旋转角度
    zoom_x = 2.0  # 设置图片相对于PDF文件在X轴上的缩放比例
    zoom_y = 2.0  # 设置图片相对于PDF文件在Y轴上的缩放比例
    trans = fitz.Matrix(zoom_x, zoom_y).preRotate(rotate)
    print("%s开始转换..." % filename)
    if doc.pageCount > 1:  # 获取PDF的页数
        for pg in tqdm(range(doc.pageCount)):
            page = doc[pg]  # 获得第pg页
            pm = page.getPixmap(matrix=trans, alpha=False)  # 将其转化为光栅文件（位数）
            new_full_name = filename.split(".")[0]  # 保证输出的文件名不变
            if not os.path.exists('./jpgs/%s' % new_full_name):
                os.mkdir('./jpgs/%s' % new_full_name)
            pm.writeImage("./jpgs/%s/%s-%s.jpg" % (new_full_name, new_full_name, pg))  # 将其输入为相应的图片格式，可以为位图，也可以为矢量图
            # 我本来想输出为jpg文件，但是在网页中都是png格式（即调用writePNG），再转换成别的图像文件前，最好查一下是否支持
    else:
        page = doc[0]
        pm = page.getPixmap(matrix=trans, alpha=False)
        new_full_name = full_path.split(".")[0]
        pm.writeImage("%s.jpg" % new_full_name)
    print("%s转换jpg完成！" % filename)
    print('耗时：', time.time() - timestamp, 's')
    print('\n\n')


# ================================================================================================================================================
    # 将分割好的jpg图片整合到ppt
    timestamp = time.time()
    print('整合jpg为ppt...')
    filename = filename.split(".")[0]
    jpg_path = './jpgs/%s' % filename
    pages = os.listdir(jpg_path)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    for index, page in enumerate(tqdm(pages)):
        # print(index)
        #Save as 'jpg' in jpgs dir
        jpg_file = "./jpgs/%s/%s-%d.jpg" % (filename,filename,index)

        #Get width/height of image
        image = Image.open(jpg_file)
        height = image.height
        width = image.width

        # #Rotate 270 degrees if horizontal
        # if height > width:
        #     adjusted = image.rotate(270, expand=True)
        #     adjusted.save(jpg_file)

        #Setup slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        left = top = 0
        slide.shapes.add_picture(jpg_file, left, top, height=prs.slide_height, width=prs.slide_width)


    prs.save('result/%s.pptx' % filename)

    print("成功保存ppt文件 %s.pptx", filename)
    print('耗时：', time.time() - timestamp, 's')
    print('\n\n')